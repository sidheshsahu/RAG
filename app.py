import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import streamlit as st
from langchain_community.document_loaders import WebBaseLoader
from bs4 import BeautifulSoup
from langchain_groq import ChatGroq
import pandas as pd
from sqlalchemy import create_engine
from langchain_community.utilities import SQLDatabase
from langchain_community.agent_toolkits import create_sql_agent
import os

# Load environment variables
load_dotenv()

# Document processing functions
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text

def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_url_text(url):
    loader = WebBaseLoader(web_paths=(url,), bs_kwargs=dict(parse_only=BeautifulSoup))
    docs = loader.load()
    text = "".join([doc.page_content for doc in docs])
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the answer is not available, respond with 'answer not available in the context'.
    
    Context:\n {context}\n
    Question:\n{question}\n
    Answer:
    """
    model = ChatGroq(model="llama3-8b-8192")
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

def query_document(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = vector_store.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain.invoke({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def query_csv_database(query):
    engine = create_engine("sqlite:///uploaded_data.db")
    db = SQLDatabase(engine=engine)
    llm = ChatGroq(model="llama3-8b-8192")
    agent_executor = create_sql_agent(llm, db=db, agent_type="zero-shot-react-description", verbose=False)
    result = agent_executor.invoke({"input": query})
    return result["output"]


# Initialize session state variables for question and answer
if "last_question" not in st.session_state:
    st.session_state.last_question = ""
if "answer" not in st.session_state:
    st.session_state.answer = ""

# Function to handle new question and answer
def handle_question(question, process_question):
    if question != st.session_state.last_question:
        st.session_state.last_question = question
        st.session_state.answer = process_question(question)
    return st.session_state.answer


# Streamlit App Layout
def main():
    st.set_page_config(page_title="Multi Agents RAG")
    st.title("Unified Document & SQL Querying App 📝🔍")

    # Tabs for Different Document Types
    tabs = st.tabs(["PDF", "DOCX", "PPTX", "URL", "CSV","XLSX"])

    with tabs[0]:  # PDF Tab
        st.header("Upload PDF Files")
        pdf_docs = st.file_uploader("Select PDF files", type="pdf", accept_multiple_files=True)
        if pdf_docs and st.button("Process PDF"):
            with st.spinner("Processing PDF..."):
                raw_text = get_pdf_text(pdf_docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("PDF processed successfully!")

    with tabs[1]:  # DOCX Tab
        st.header("Upload DOCX Files")
        docx_docs = st.file_uploader("Select DOCX files", type="docx", accept_multiple_files=True)
        if docx_docs and st.button("Process DOCX"):
            with st.spinner("Processing DOCX..."):
                raw_text = get_docx_text(docx_docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("DOCX processed successfully!")

    with tabs[2]:  # PPTX Tab
        st.header("Upload PPTX Files")
        pptx_docs = st.file_uploader("Select PPTX files", type="pptx", accept_multiple_files=True)
        if pptx_docs and st.button("Process PPTX"):
            with st.spinner("Processing PPTX..."):
                raw_text = get_pptx_text(pptx_docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("PPTX processed successfully!")

    with tabs[3]:  # URL Tab
        st.header("Enter a URL")
        url_input = st.text_input("Enter a URL to extract text content from")
        if url_input and st.button("Process URL"):
            with st.spinner("Processing URL..."):
                raw_text = get_url_text(url_input)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("URL processed successfully!")

    with tabs[4]:  # CSV Tab
        st.header("Upload CSV File")
        csv_file = st.file_uploader("Select a CSV file", type="csv")
        if csv_file:
            df = pd.read_csv(csv_file)
            st.write("Data Preview:", df.head())
            engine = create_engine("sqlite:///uploaded_data.db")
            df.to_sql("uploaded_table", engine, index=False, if_exists="replace")
            st.success("CSV loaded into SQLite database.")
            query = st.text_input("Ask a question about your CSV data")
            if query and st.button("Run Query"):
                with st.spinner("Running SQL Query..."):
                    try:
                        result = query_csv_database(query)
                        st.write("Query Result:", result)
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
    
    with tabs[5]:  # Excel Tab
        st.header("Upload Excel File")
        excel_file = st.file_uploader("Select an Excel file", type="xlsx")
        if excel_file:
            df = pd.read_excel(excel_file)
            st.write("Data Preview:", df.head())
            engine = create_engine("sqlite:///uploaded_data.db")
            df.to_sql("uploaded_table", engine, index=False, if_exists="replace")
            st.success("Excel loaded into SQLite database.")
            query = st.text_input("Ask a question about your Excel data")
            if query and st.button("Run Excel Query"):
                with st.spinner("Running SQL Query..."):
                    try:
                        result = query_csv_database(query)
                        st.write("Query Result:", result)
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
    
    
    
    
    # Question-Answering Section
    st.subheader("Question Answering")
    user_question = st.text_input("Ask a question about the processed documents (PDF, DOCX, PPTX, URL)")
    if user_question:
        answer = query_document(user_question)
        st.write("Answer:", answer)

if __name__ == "__main__":
    main()
















# import streamlit as st
# import pandas as pd
# from sqlalchemy import create_engine
# from PyPDF2 import PdfReader
# from docx import Document
# from pptx import Presentation
# from bs4 import BeautifulSoup
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_google_genai import GoogleGenerativeAIEmbeddings
# from langchain_community.vectorstores import FAISS
# from langchain_google_genai import ChatGoogleGenerativeAI
# from langchain.chains.question_answering import load_qa_chain
# from langchain.prompts import PromptTemplate
# from langchain_groq import ChatGroq
# from langchain_community.utilities import SQLDatabase
# from langchain_community.agent_toolkits import create_sql_agent
# from langchain_community.document_loaders import WebBaseLoader
# from dotenv import load_dotenv

# # Load environment variables
# load_dotenv()

# # SQLite database engine setup
# engine = create_engine("sqlite:///uploaded_data.db")

# # ===== Document Processing Functions ===== #
# def get_pdf_text(pdf_docs):
#     text = ""
#     for pdf in pdf_docs:
#         reader = PdfReader(pdf)
#         for page in reader.pages:
#             text += page.extract_text()
#     return text

# def get_docx_text(docx_docs):
#     text = ""
#     for docx in docx_docs:
#         doc = Document(docx)
#         for para in doc.paragraphs:
#             text += para.text + "\n"
#     return text

# def get_pptx_text(pptx_docs):
#     text = ""
#     for pptx in pptx_docs:
#         presentation = Presentation(pptx)
#         for slide in presentation.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + "\n"
#     return text

# def get_url_text(url):
#     loader = WebBaseLoader(web_paths=(url,), bs_kwargs=dict(parse_only=BeautifulSoup))
#     docs = loader.load()
#     return "".join([doc.page_content for doc in docs])

# def get_text_chunks(text):
#     splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
#     return splitter.split_text(text)

# def get_vector_store(text_chunks):
#     embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
#     vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
#     vector_store.save_local("faiss_index")

# def query_document(question):
#     embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
#     vector_store = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
#     docs = vector_store.similarity_search(question)
#     chain = get_conversational_chain()
#     response = chain.invoke({"input_documents": docs, "question": question}, return_only_outputs=True)
#     return response["output_text"]

# def get_conversational_chain():
#     prompt = PromptTemplate(
#         template="""
#         Answer the question as detailed as possible from the provided context. 
#         If the answer is not available, respond with 'answer not available in the context'.

#         Context:\n {context}\n
#         Question:\n{question}\n
#         Answer:
#         """, input_variables=["context", "question"]
#     )
#     model = ChatGroq(model="llama3-8b-8192")
#     return load_qa_chain(model, chain_type="stuff", prompt=prompt)

# # ===== SQL Query Function ===== #
# def query_database(query):
#     result = ""
#     try:
#         with engine.connect() as conn:
#             result = conn.execute(query).fetchall()
#     except Exception as e:
#         result = f"An error occurred: {e}"
#     return result

# def query_csv_database(user_query):
#     try:
#         # Connect to the SQLite database
#         engine = create_engine("sqlite:///uploaded_data.db")
#         with engine.connect() as conn:
#             # Execute the user's query
#             result = conn.execute(text(user_query))
#             # Fetch all results as a DataFrame
#             df_result = pd.DataFrame(result.fetchall(), columns=result.keys())
#         return df_result
#     except Exception as e:
#         raise ValueError(f"SQL query failed: {str(e)}")

# # Main function to handle app flow
# def main():
#     st.set_page_config(page_title="Unified Querying App")
#     st.title("Unified Document & SQL Querying App 📝🔍")

#     # Flags to track data uploads
#     data_loaded = False
#     is_csv_or_excel = False

#     # Uploading CSV and Excel Files
#     csv_file = st.file_uploader("Select a CSV file", type="csv")
#     if csv_file:
#         df = pd.read_csv(csv_file)
#         st.write("Data Preview:", df.head())
#         engine = create_engine("sqlite:///uploaded_data.db")
#         df.to_sql("uploaded_table", engine, index=False, if_exists="replace")
#         st.success("CSV loaded into SQLite database.")
#         data_loaded = True
#         is_csv_or_excel = True

#     excel_file = st.file_uploader("Select an Excel file", type="xlsx")
#     if excel_file:
#         df = pd.read_excel(excel_file)
#         st.write("Data Preview:", df.head())
#         engine = create_engine("sqlite:///uploaded_data.db")
#         df.to_sql("uploaded_table", engine, index=False, if_exists="replace")
#         st.success("Excel loaded into SQLite database.")
#         data_loaded = True
#         is_csv_or_excel = True

#     # Unified Search Bar
#     st.subheader("Ask a Question 📋")
#     user_query = st.text_input("Ask a question about your uploaded data (CSV, Excel, PDF, DOCX, PPTX, URL)")

#     if user_query:
#         with st.spinner("Processing your query..."):
#             if is_csv_or_excel:
#                 # Query the SQLite database for CSV/Excel data
#                 try:
#                     result = query_csv_database(user_query)
#                     st.write("Query Result:", result)
#                 except ValueError as e:
#                     st.error(f"Error: {e}")
#             else:
#                 st.warning("Please upload a document (PDF, DOCX, PPTX, or URL) for search queries.")

# if __name__ == "__main__":
#     main()